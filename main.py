import itertools
import os

import pygame
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from pptx.util import Pt

import Ada_pygame_prototype


def generate_slide_trials():
    pr1 = Presentation()

    # Refers to the layout of each slide.
    # 0 = Title slide, 1 = Title and Content, 3 = Section Header.
    slide1_register = pr1.slide_layouts[0]

    # Slide 1:
    # Add initial slide to presentation
    slide1 = pr1.slides.add_slide(slide1_register)

    # Add text to layouts.
    title = slide1.shapes.title
    # Placeholder = Item in the layout.
    subtitle = slide1.placeholders[1]

    # Insert text into the placeholders.
    title.text = "Test xixixi"
    subtitle.text = "lah lah lah"

    # Slide 2:
    slide2_register = pr1.slide_layouts[1]
    slide2 = pr1.slides.add_slide(slide2_register)
    # Edit bullet point slide.
    title2 = slide2.shapes.title
    title2.text = "Now for some bullet points"

    bullet_point_box = slide2.shapes
    bullet_points_lvl1 = bullet_point_box.placeholders[1]
    bullet_points_lvl1.text = "lah lah lah hei hei hei"

    # Create the second level of bullet points.
    bullet_points_lvl2 = bullet_points_lvl1.text_frame.add_paragraph()
    bullet_points_lvl2.text = "du du du"
    bullet_points_lvl2.level = 1

    bullet_points_lvl3 = bullet_points_lvl1.text_frame.add_paragraph()
    bullet_points_lvl3.text = "du du du 2"
    bullet_points_lvl3.level = 2

    bullet_points_lvl4 = bullet_points_lvl1.text_frame.add_paragraph()
    bullet_points_lvl4.text = "du du du 3"
    bullet_points_lvl4.level = 3

    # Slide 3:
    slide3_register = pr1.slide_layouts[5]
    slide3 = pr1.slides.add_slide(slide3_register)

    title3 = slide3.shapes.title
    title3.text = "Picture"

    # Add image
    img1 = "image_1.jpg"

    from_left = Inches(3)
    from_top = Inches(2)
    add_picture = slide3.shapes.add_picture(img1, from_left, from_top)

    # Slide 4
    slide4_register = pr1.slide_layouts[5]
    slide4 = pr1.slides.add_slide(slide4_register)

    title4 = slide4.shapes.title
    title4.text = "Shapework"

    # Create shapes
    # Shape 1
    left1 = top1 = width1 = height1 = Inches(2)
    add_shape1 = slide4.shapes.add_shape(MSO_SHAPE_TYPE.PICTURE, left1, top1, width1, height1)

    # Shape 2
    left2 = Inches(6)
    top2 = Inches(2)
    width2 = height2 = Inches(2)
    shape2 = slide4.shapes.add_shape(MSO_SHAPE_TYPE.PICTURE, left2, top2, width2, height2)

    # Change the shape's main color
    fill_shape2 = shape2.fill
    fill_shape2.solid()
    fill_shape2.fore_color.theme_color = MSO_THEME_COLOR_INDEX.LIGHT_2

    # Slide 4:
    slide5_register = pr1.slide_layouts[5]
    slide5 = pr1.slides.add_slide(slide5_register)

    # Title:
    title5 = slide5.shapes.title
    title5.text = "Graphs"

    # Build graph
    graph_info = CategoryChartData()
    graph_info.categories = ["A", "B", "C"]
    graph_info.add_series("Series 1", (15, 11, 18))

    # Add graph to slide with positioning
    left_graph = Inches(2)
    top_graph = Inches(2)
    width_graph = Inches(6)
    height_graph = Inches(4)
    # slide5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    #                         left_graph, top_graph, width_graph, height_graph,
    #                         graph_info)
    graph1_frame = slide5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                           left_graph, top_graph, width_graph, height_graph,
                                           graph_info)
    graph1 = graph1_frame.chart

    # Edit graph
    category_axis = graph1.category_axis
    category_axis.has_major_gridlines = True
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(24)

    # Slide 6:
    slide6_register = pr1.slide_layouts[5]
    slide6 = pr1.slides.add_slide(slide6_register)

    title6 = slide6.shapes.title
    title6.text = "Table"

    # Add Graph to slide
    left_table = Inches(1)
    top_table = Inches(2)
    width_table = Inches(8)
    height_table = Inches(3)
    # slide6.shapes.add_table(3,4, left_table, top_table,
    #                         width_table, height_table)
    table1_frame = slide6.shapes.add_table(3, 4, left_table, top_table,
                                           width_table, height_table)

    # Populating a table
    table1 = table1_frame.table
    cell = table1.cell(0, 0)
    cell.text = "Insert the title here"

    cell = table1.cell(1, 0)
    cell.text = "Insert Value here"

    # Rotate the shape
    shape2.rotation = 90

    # Slide 7
    slide7_register = pr1.slide_layouts[0]
    slide7 = pr1.slides.add_slide(slide7_register)

    title7 = slide7.shapes.title
    title7.text = "Hyperlink"

    # Adding hyperlinks
    para1 = slide7.placeholders[1].text_frame.paragraphs[0]
    addrun1 = para1.add_run()
    addrun1.text = "Google Hyperlink"
    hlink1 = addrun1.hyperlink
    hlink1.address = "https://www.google.com"

    pr1.save("trial_slides.pptx")


def generate_pygame_window_trials():
    # Activate, initiate pygame and give permission to use pygame's functionality.
    pygame.init()

    # Define the RGB value for the color
    zima_green = (73, 232, 56)
    dudu_purple = (102, 0, 204)
    background_black = (8, 8, 8)

    # Assign values to X and Y variable
    width_window = 400
    height_window = 400

    # Create the display surface object of specific dimension (width, height)
    display_surface = pygame.display.set_mode((width_window, height_window))

    # Set the pygame window name
    pygame.display.set_caption('Text display')

    # Create a font object, the 1st parameter is the font file, and 2nd parameter is size of the font
    font = pygame.font.Font('freesansbold.ttf', 32)

    # Create a text surface object, on which text is drawn on it
    text = font.render('Geeks for Geeks', True, zima_green, dudu_purple)

    # Create a rectangular object for the text surface object
    text_rectangular = text.get_rect()

    # Set the center of the rectangular object
    text_rectangular.center = (width_window // 2, height_window // 2)

    # Infinite loop
    while True:
        # Completely fill the surface object with the background black
        display_surface.fill(background_black)

        # Copy the surface object to the display surface object at the center coordinate
        display_surface.blit(text, text_rectangular)

        # Iterate over the list of Event objects that was returned by pygame.event.get() method.
        for event in pygame.event.get():
            # If event object type is QUIT then quitting the pygame and program both.
            if event.type == pygame.QUIT:
                # Deactivates the pygame library
                pygame.quit()
                # Quit the program.
                quit()

            # Draws the surface object to the screen.
            pygame.display.update()


# A more sophisticated pygame tutorial
# https://www.reddit.com/r/pygame/comments/d175oj/how_do_i_display_text_for_a_set_amount_of_time/
# Some instructions from the expert:
# First let break it down.
# Text need to keep a time with it. Text will need to be in a container holding data.
# Display would need a container. When times up then remove it.

# Things to avoid.
# time.sleep. Never use it in pygame mainloop.
# Don't use function to draw text directly to main surface.

# Your function draw text function. Need to be rewritten.
# You should not be creating font every time you want a new text.
# It should return text surface or text surface in a container with other data.


class Scene:
    def on_draw(self, surface): pass

    def on_update(self, delta): pass

    def on_event(self, event): pass


class Manager:
    """
    Create the canvas/window that contains surfaces.
    Run the main process of the game.
    """

    @classmethod
    def create(cls, title, width, height, center=False):
        if center:
            os.environ['SDL_VIDEO_CENTERED'] = '1'

        # Basic pygame setup
        pygame.display.set_caption(title)
        cls.surface = pygame.display.set_mode((width, height))
        cls.rect = cls.surface.get_rect()
        cls.clock = pygame.time.Clock()
        cls.running = False
        cls.delta = 0
        cls.fps = 60

        cls.scene = Scene()

        cls.is_text_shown = True
        # cls.is_task_shown = False

        # Deisplayed text shown:
        cls.time_threshold = 3000
        cls.timer = 0

        text = "wai bi ba bo?"
        font_color = (73, 232, 56)
        anchor = "topleft"
        position = (20, 0)
        font_text = pygame.font.Font(None, 24)  # TODO: make the font a tunable parameter.
        cls.image_text = font_text.render(text, 1, font_color)
        cls.rect_text = cls.image_text.get_rect()
        setattr(cls.rect_text, anchor, position)

        task_text = "Wubba lubba dub dub!"
        font_color_task = pygame.Color("dodgerblue")
        position_task = (50, 60)
        font_task = pygame.font.Font(None, 48)
        cls.image_task = font_task.render(task_text, 1, font_color_task)
        cls.rect_task = cls.image_task.get_rect()
        setattr(cls.rect_task, anchor, position_task)

    @classmethod
    def mainloop(cls):
        cls.running = True
        while cls.running:
            for event in pygame.event.get():
                # Used a loop because has to listen to the keyboard pressed event.
                # Add text to the text list when necessary (space key is pressed).
                cls.scene.on_event(event)

            # Refresh and update the text list according to the timers,
            # eliminate the overtime ones and adjust placements.
            cls.scene.on_update(cls.delta)
            # Display texts.
            # cls.scene.on_draw(cls.surface)    # TODO: unveil this will cause override.

            pygame.display.flip()
            cls.delta = cls.clock.tick(cls.fps)  # The timer to monitor elapsed time: 16/17ms each time.
            # tick():This method should be called once per frame.
            # It will compute how many milliseconds have passed since the previous call.
            # If you pass the optional framerate argument the function will delay to keep the game running slower
            # than the given ticks per second. For example if we pass 10 as argument
            # the program will never run at more than 10 frames per second.

            # This part's reference:
            # https://stackoverflow.com/questions/44206901/how-to-display-text-for-only-a-certain-amount-of-time
            elapsed = cls.delta
            if cls.is_text_shown:
                # Display the text:

                cls.surface.blit(cls.image_text, cls.rect_text)

                cls.timer = cls.timer + elapsed  # Timer's unit is "ms".
                if cls.timer > cls.time_threshold:
                    cls.is_text_shown = False
                    cls.timer = 0
                    # TODO: erase the shown text
                    cls.surface.fill("black")
            elif cls.is_text_shown is False:
                # Display the task that aims to erase users short term memory

                cls.surface.blit(cls.image_task, cls.rect_task)

                cls.timer = cls.timer + elapsed  # Timer's unit is "ms".
                if cls.timer > cls.time_threshold:
                    cls.is_text_shown = True
                    cls.timer = 0
                    cls.surface.fill("black")


class TextTimed:
    # Produce a single timed text (the text that is shown for a given amount of time).
    def __init__(self, font, text, foreground, position, timed=3000, anchor="topleft"):
        # TODO: change the timed into a tunable parameter.
        self.image = font.render(text, 1, foreground)  # Render: draw text on a new surface.
        # Syntax: render(text, antialias, color, background=None) -> Surface
        self.rect = self.image.get_rect()
        setattr(self.rect, anchor, position)  # Set the value of the specified attribute of the specified object.
        # Syntax: setattr(object, attribute, value)
        # Here the "topleft" is an attribute of the rectangular, which means a tuple of two integers: (left, top)

        self.timed = timed  # Duration

    def draw(self, surface):
        # When it is still within the text display duration, keep showing texts
        if self.timed > 0:
            surface.blit(self.image, self.rect)  # Syntax: blit(source, dest, area=None, special_flags=0) -> Rect
            # Source – Draws a source Surface onto this Surface
            # dest – The draw can be positioned with the dest argument.
            # area -A Rect can also be passed as the destination
            # and the topleft corner of the rectangle will be used as the position for the blit

    # Count down the time for displaying texts.
    def update(self, delta):
        self.timed -= delta


class Example(Scene):  # Class-Example is inherited from the class-Scene.
    def __init__(self):
        self.font = pygame.font.Font(None, 24)
        self.timed_text = []
        self.count = itertools.count(start=1, step=1)  # An infinite iterator

    def position_text(self):
        y = itertools.count(Manager.rect.bottom - 50, -30)  # Move the old texts up
        for text in self.timed_text[::-1]:  # Present reversely, the old ones are stacked up to higher places
            text.rect.y = next(y)  # Returns the next item in an iterator.

    # Display the update set of texts.
    def on_draw(self, surface):
        surface.fill(pygame.Color("white"))
        for text in self.timed_text:
            text.draw(surface)

    # Update the set of texts that are shown on the canvas, delete the overtime ones, then reposition all of them.
    def on_update(self, delta):
        timed_text = []
        for text in self.timed_text:
            # Each text has a timer.
            text.update(delta)
            if text.timed > 0:
                timed_text.append(text)

        self.timed_text = timed_text
        self.position_text()

    def on_event(self, event):
        if event.type == pygame.QUIT:
            Manager.running = False
        elif event.type == pygame.KEYDOWN:
            if event.key == pygame.K_SPACE:
                text = "Timed Text {}".format(next(self.count))
                # color = pygame.Color("dodgerblue")
                # TODO: change this color to a tunable variable.
                color = (73, 232, 56)  # A light green that has higher contracts for OHMDs. Cool, it works.
                timed = TextTimed(self.font, text, color, (20, 0))
                self.timed_text.append(timed)  # Update the texts cluster.
                self.position_text()


class ExampleAttentionShift(Scene):  # Class-Example is inherited from the class-Scene.
    def __init__(self):
        self.font = pygame.font.Font(None, 24)  # TODO: make the font a tunable parameter.
        self.timed_text = []
        self.count = itertools.count(start=1, step=1)  # An infinite iterator

        # Initiate the text registrations:
        self.buffer_texts = []
        self.buffer_text_contents = ['a', 'b', 'c', 'd', 'e', 'f', 'g']
        # TODO: This stores a large amount of text, waiting to be broke into chunks and dispalyed.
        # TODO: change this color to a tunable variable.
        # TODO: insert a paragraph split function here.
        color = (73, 232, 56)  # A light green that has higher contracts for OHMDs. Cool, it works.
        for text_content in self.buffer_text_contents:
            instanced_text_content = TextTimed(self.font, text_content, color, (20, 0))
            self.buffer_texts.append(instanced_text_content)

    def position_text(self):
        y = itertools.count(Manager.rect.bottom - 50, -30)  # Move the old texts up
        for text in self.timed_text[::-1]:  # Present reversely, the old ones are stacked up to higher places
            text.rect.y = next(y)  # Returns the next item in an iterator.

    # Display the update set of texts.
    def on_draw(self, surface):
        surface.fill(pygame.Color("black"))  # TODO: make the background color a tunable parameter.
        for text in self.timed_text:
            text.draw(surface)
        # self.timed_text[0].draw(surface)

    # Update the set of texts that are shown on the canvas, delete the overtime ones, then reposition all of them.
    def on_update(self, delta):
        timed_text = []
        for text in self.timed_text:
            # Each text has a timer.
            text.update(delta)
            if text.timed > 0:
                timed_text.append(text)

        self.timed_text = timed_text
        self.position_text()

    def on_event(self, event):
        if event.type == pygame.QUIT:
            Manager.running = False
        elif event.type == pygame.KEYDOWN:
            if event.key == pygame.K_SPACE:
                text = "Timed Text {}".format(next(self.count))
                # color = pygame.Color("dodgerblue")
                # TODO: change this color to a tunable variable.
                color = (73, 232, 56)  # A light green that has higher contracts for OHMDs. Cool, it works.
                timed = TextTimed(self.font, text, color, (20, 0))
                self.timed_text.append(timed)  # Update the texts cluster.
                self.position_text()


def example_main():
    pygame.init()
    Manager.create("Example Timed Text", 800, 600, True)
    Manager.scene = Example()
    Manager.mainloop()


# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    # generate_slide_trials()

    # ada_generator = PptxGenerate(start_pos="on_OHMD", time_off=5, time_on=5, amount_text=50,
    #                              num_shifts=10, subtask_type="default", name_ppt="trial.pptx",
    #                              text="lalalalalalalala")
    # ada_generator.generate_slides()

    # Manage the slides transition time, .
    # Aspose will produce water marker, removing it seems to require a licensed account.
    # ada_generator.set_transition_time()

    # generate_pygame_window_trials()
    # example_main()

    Ada_pygame_prototype.run_prototype()
