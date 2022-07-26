from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt
import aspose.slides as slides


class PptxGenerate:
    def __init__(self, start_pos, time_off, time_on, amount_text,
                 num_shifts, subtask_type, name_ppt, text):
        # Initialize the with-in class global parameters from arguments
        self.starting_slide = start_pos
        self.duration_off_OHMD = time_off
        self.duration_on_OHMD = time_on
        self.amount_text_OHMD = amount_text
        self.number_slides = num_shifts + 1
        self.task_type_off_OHMD = subtask_type
        self.name_ppt = name_ppt
        self.text = text

        # Initialize other global variables
        # Referenced documents: https://python-pptx.readthedocs.io/en/latest/user/presentations.html
        self.pr = Presentation()
        # Register slides according to the number of attentional shifts
        self.slides_registers = []
        self.slides = []

        # Initiate the constant variables
        self.from_left_textbox = Inches(5)
        self.from_top_textbox = Inches(3)
        self.width_textbox = Inches(5)
        self.height_textbox = Inches(5)

        self.from_left_task = Inches(1)
        self.from_top_task = Inches(1)

    def generate_slides(self):
        # Configurations
        # Register and create slides
        for i in range(self.number_slides):
            slide_register = self.pr.slide_layouts[6]  # Index = 6 means all blank slide.
            slide = self.pr.slides.add_slide(slide_register)
            self.slides_registers.append(slide_register)
            self.slides.append(slide)

        # Determine the starting slide
        # Blogs referenced: https://stackoverflow.com/questions/24729098/text-color-in-python-pptx-module.
        if self.starting_slide == "on_OHMD":
            for i in range(len(self.slides)):
                if i % 2 == 0:    # The odd slides
                    # # Create text frames
                    # para = self.slides[i].placeholders[1].text_frame.paragraphs[0]
                    # para.text = "ha ha ha xi xi xi"
                    # Create textboxes
                    textbox = self.slides[i].shapes.add_textbox(self.from_left_textbox,
                                                              self.from_top_task,
                                                              self.width_textbox,
                                                              self.height_textbox)
                    tf = textbox.text_frame

                    p = tf.paragraphs[0]     # Defaultly there is one paragraph.
                    p.text = "This is a first paragraph"
                    p.font.size = Pt(30)
                    # p.font.bold = True
                    # p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                    p.font.color.rgb = RGBColor(73, 232, 56)

                    # p = tf.add_paragraph()    # For the second paragraph, needs to be added
                    # p.text = "This is a second paragraph"
                    # p.font.size = Pt(11)
        elif self.starting_slide == "off_OHMD":
            pass

        # Save the slide
        self.pr.save(self.name_ppt)

    def set_transition_time(self):
        # Reference: https://gist.github.com/aspose-com-gists/45c4d8504096fb877c50c8289130cbf1
        # Instantiate Presentation class that represents a presentation file
        with slides.Presentation(self.name_ppt) as pres:
            # Apply circle type transition on slide 1
            pres.slides[0].slide_show_transition.type = slides.slideshow.TransitionType.CIRCLE

            # Set the transition time of 3 seconds
            pres.slides[0].slide_show_transition.advance_on_click = True
            pres.slides[0].slide_show_transition.advance_after_time = 3000

            # Apply comb type transition on slide 2
            pres.slides[1].slide_show_transition.type = slides.slideshow.TransitionType.COMB

            # Set the transition time of 5 seconds
            pres.slides[1].slide_show_transition.advance_on_click = True
            pres.slides[1].slide_show_transition.advance_after_time = 5000

            # Apply zoom type transition on slide 3
            pres.slides[2].slide_show_transition.type = slides.slideshow.TransitionType.ZOOM

            # Set the transition time of 7 seconds
            pres.slides[2].slide_show_transition.advance_on_click = True
            pres.slides[2].slide_show_transition.advance_after_time = 7000

            # Write the presentation to disk
            pres.save("trial_transition.pptx", slides.export.SaveFormat.PPTX)
